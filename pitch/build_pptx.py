#!/usr/bin/env python3
"""
Builds OXAR-pitch.pptx.

This file is the deck. Not a copy of it, not an export of it — the deck. Which is
the whole reason it exists: the previous deck was React, and a React deck cannot be
opened by the person who needs to change one word before a meeting. A .pptx opens in
Google Slides, Keynote and PowerPoint, and every figure in it still comes from one
place that can be reviewed in a diff.

Every number here was checked on 2026-08-16 against the production database, the
code, or a published source named on the slide. Nothing goes in that we cannot
stand behind.

    python3 pitch/build_pptx.py [output.pptx]

LAYOUT NOTE, learned the hard way twice: python-pptx cannot measure text, so a
layout that stacks blocks by estimating the height of the one above will overlap
it the first time a headline wraps one line further than expected. Fixed slots
were the second attempt and only moved the problem — a slot generous enough for
the longest string wastes the space on every other slide.

So we measure. `textfit` reads the real DM Sans and does the renderer's own greedy
wrap, which means line counts are exact rather than guessed. Every block flows
under the measured bottom of the one above it, and `_check` asserts at build time
that nothing overlaps, overflows, or crosses into a picture.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import textfit  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
ART = ROOT / "web" / "public" / "pitch" / "collage"

# --- The look -----------------------------------------------------------------
# Same three colours and one typeface as the app and the landing: black, white, and
# the violet that only ever marks the stressed word.

FONT = "DM Sans"
INK_DARK, PAPER_DARK = RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0x00, 0x00, 0x00)
INK_LIGHT, PAPER_LIGHT = RGBColor(0x0A, 0x0A, 0x0A), RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
MUTED_DARK, MUTED_LIGHT = RGBColor(0x8C, 0x8C, 0x8C), RGBColor(0x6B, 0x6B, 0x6B)
FAINT_DARK, FAINT_LIGHT = RGBColor(0x5E, 0x5E, 0x5E), RGBColor(0x9A, 0x9A, 0x9A)
RULE_DARK, RULE_LIGHT = RGBColor(0x28, 0x28, 0x28), RGBColor(0xE2, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)
PAD = Inches(0.85)

# Where a slide starts, and the air between blocks. Everything else is measured.
Y_TOP = Inches(0.95)
Y_TOP_STATEMENT = Inches(2.05)
GAP_KICKER, GAP_TITLE, GAP_SUB = Inches(0.30), Inches(0.34), Inches(0.52)
Y_FOOT = H - Inches(1.10)

# Where the picture starts, and how much room that leaves the words. These two are a
# pair: widen one and the other has to give, or text slides under the photograph.
PIC_L, PIC_L_WIDE = Inches(8.3), Inches(7.0)
COL_FULL = W - 2 * PAD
COL_SPLIT = PIC_L - PAD - Inches(0.35)
COL_TITLE = PIC_L_WIDE - PAD - Inches(0.35)


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H

    # -- primitives ------------------------------------------------------------

    def _slide(self, light: bool, image: str | None, *, wide: bool = False):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = PAPER_LIGHT if light else PAPER_DARK
        # Picture first, so every word lands on top of it rather than under it.
        if image:
            self._picture(s, image, wide=wide)
        return s

    def _picture(self, s, name: str, *, wide: bool) -> None:
        path = ART / name
        if not path.exists():
            raise FileNotFoundError(f"missing deck art: {path}")
        from PIL import Image as PILImage

        iw, ih = PILImage.open(path).size
        # The right band, inset from the edges so nothing looks pasted on.
        box_l = PIC_L_WIDE if wide else PIC_L
        box_w = W - box_l - Inches(0.5)
        box_t, box_h = Inches(0.7), H - Inches(1.4)
        scale = min(box_w / iw, box_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        s.shapes.add_picture(
            str(path),
            Emu(int(box_l + (box_w - w) / 2)),
            Emu(int(box_t + (box_h - h) / 2)),
            Emu(w), Emu(h),
        )

    def _box(self, s, left, top, width, height):
        tb = s.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        return tf

    def _write(self, s, left, top, width, segments, *, size, color, bold=False, spacing=1.0):
        """Writes a paragraph and returns how tall it actually came out. Built from
        (text, is_accent) pairs, so one word carries the violet italic without the
        whole line changing typeface."""
        plain = "".join(t for t, _ in segments)
        height = textfit.height_emu(plain, width, size, bold=bold, spacing=spacing)
        tf = self._box(s, left, top, width, height)
        p = tf.paragraphs[0]
        p.line_spacing = spacing
        for text, accent in segments:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, bool(accent)
            f.color.rgb = ACCENT if accent else color
        return height

    def _rule(self, s, top, width, light: bool):
        line = s.shapes.add_connector(1, PAD, top, PAD + width, top)
        line.line.color.rgb = RULE_LIGHT if light else RULE_DARK
        line.line.width = Pt(0.75)

    # -- slide furniture -------------------------------------------------------

    def _head(self, s, kicker, title, sub, light, col, *, title_size=40, top=Y_TOP):
        """Kicker, headline, subtitle — each under the measured bottom of the last.
        Returns where the body may start."""
        ink = INK_LIGHT if light else INK_DARK
        muted = MUTED_LIGHT if light else MUTED_DARK

        y = top
        y += self._write(s, PAD, y, col, [(f"[ {kicker} ]", False)], size=13, color=muted)
        y += GAP_KICKER
        y += self._write(s, PAD, y, col, title, size=title_size, color=ink, bold=True, spacing=0.95)
        y += GAP_TITLE
        if sub:
            y += self._write(s, PAD, y, min(col, Inches(6.6)), [(sub, False)],
                             size=14.5, color=muted, spacing=1.35)
            y += GAP_SUB
        return Emu(int(y))

    def _foot(self, s, text, light, col=None):
        if not text:
            return
        width = col or COL_FULL
        h = textfit.height_emu(text, width, 10.5, spacing=1.3)
        self._write(s, PAD, Emu(int(H - Inches(0.6) - h)), width, [(text, False)],
                    size=10.5, color=FAINT_LIGHT if light else FAINT_DARK, spacing=1.3)

    # -- slide kinds -----------------------------------------------------------

    def title(self, title, sub, image, note):
        s = self._slide(False, image, wide=True)
        y = Inches(2.45)
        y += self._write(s, PAD, y, COL_TITLE, title, size=48, color=INK_DARK, bold=True, spacing=0.95)
        y += Inches(0.45)
        self._write(s, PAD, y, COL_TITLE, [(sub, False)], size=15, color=MUTED_DARK, spacing=1.4)
        s.notes_slide.notes_text_frame.text = note
        return s

    def statement(self, kicker, title, sub, light, image, note, footer=None):
        s = self._slide(light, image)
        col = COL_SPLIT if image else COL_FULL
        self._head(s, kicker, title, sub, light, col, title_size=44, top=Y_TOP_STATEMENT)
        self._foot(s, footer, light, col)
        s.notes_slide.notes_text_frame.text = note
        return s

    def columns(self, kicker, title, cols, light, image, note, footer=None):
        s = self._slide(light, image)
        col = COL_SPLIT if image else COL_FULL
        y0 = self._head(s, kicker, title, None, light, col)
        ink = INK_LIGHT if light else INK_DARK
        muted = MUTED_LIGHT if light else MUTED_DARK

        n = len(cols)
        gutter = Inches(0.45)
        width = Emu(int((col - gutter * (n - 1)) / n))
        y = y0
        self._rule(s, y - Inches(0.28), col, light)
        for i, (label, body) in enumerate(cols):
            x = Emu(int(PAD + i * (width + gutter)))
            self._write(s, x, y, width, [(label, False)], size=14, color=ink)
            self._write(s, x, y + Inches(0.48), width, [(body, False)], size=11.5, color=muted, spacing=1.4)
        self._foot(s, footer, light, col)
        s.notes_slide.notes_text_frame.text = note
        return s

    def stats(self, kicker, title, sub, items, light, image, note, footer=None):
        s = self._slide(light, image)
        col = COL_SPLIT if image else COL_FULL
        y0 = self._head(s, kicker, title, sub, light, col)
        ink = INK_LIGHT if light else INK_DARK
        muted = MUTED_LIGHT if light else MUTED_DARK
        faint = FAINT_LIGHT if light else FAINT_DARK

        n = len(items)
        gutter = Inches(0.4)
        width = Emu(int((col - gutter * (n - 1)) / n))
        y = y0
        for i, (figure, label, note_text) in enumerate(items):
            x = Emu(int(PAD + i * (width + gutter)))
            self._write(s, x, y, width, [(figure, False)], size=36, color=ink, bold=True, spacing=0.9)
            self._write(s, x, y + Inches(0.66), width, [(label, False)], size=11.5, color=muted)
            self._write(s, x, y + Inches(1.0), width, [(note_text, False)], size=10, color=faint, spacing=1.3)
        self._foot(s, footer, light, col)
        s.notes_slide.notes_text_frame.text = note
        return s

    def rows(self, kicker, title, items, light, image, note, footer=None):
        s = self._slide(light, image)
        col = COL_SPLIT if image else COL_FULL
        y = self._head(s, kicker, title, None, light, col)
        ink = INK_LIGHT if light else INK_DARK
        muted = MUTED_LIGHT if light else MUTED_DARK

        step = Emu(int((Y_FOOT - Inches(0.25) - y) / len(items)))
        label_w = Inches(2.2)
        for label, body, strong in items:
            self._rule(s, y, col, light)
            self._write(s, PAD, y + Inches(0.16), label_w, [(label, False)], size=12.5, color=ink if strong else muted, bold=strong)
            self._write(s, PAD + label_w + Inches(0.4), y + Inches(0.16),
                        col - label_w - Inches(0.4),
                        [(body, False)], size=11.5, color=ink if strong else muted, spacing=1.32)
            y = Emu(int(y + step))
        self._foot(s, footer, light, col)
        s.notes_slide.notes_text_frame.text = note
        return s

    def steps(self, kicker, title, items, light, image, note, footer=None):
        s = self._slide(light, image)
        col = COL_SPLIT if image else COL_FULL
        y = self._head(s, kicker, title, None, light, col)
        ink = INK_LIGHT if light else INK_DARK
        muted = MUTED_LIGHT if light else MUTED_DARK

        step = Inches(1.1)
        for i, body in enumerate(items, start=1):
            self._rule(s, y, col, light)
            self._write(s, PAD, y + Inches(0.2), Inches(1.1), [(f"0{i}", False)], size=26, color=ink, bold=True)
            self._write(s, PAD + Inches(1.35), y + Inches(0.3), col - Inches(1.35),
                        [(body, False)], size=12.5, color=muted, spacing=1.32)
            y += step
        self._foot(s, footer, light, col)
        s.notes_slide.notes_text_frame.text = note
        return s

    def check(self) -> list[str]:
        """Every text block against every other, plus the frame and the pictures.
        Heights are measured, not assumed, so a complaint here is a real one — and
        silence means the file is actually clean rather than merely built."""
        faults: list[str] = []
        for i, s in enumerate(self.prs.slides, start=1):
            texts = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
            pics = [sh for sh in s.shapes if sh.shape_type == 13]
            for a in range(len(texts)):
                A = texts[a]
                for b in range(a + 1, len(texts)):
                    B = texts[b]
                    if (A.left < B.left + B.width and B.left < A.left + A.width
                            and A.top < B.top + B.height and B.top < A.top + A.height):
                        faults.append(f"slide {i}: “{A.text_frame.text[:22]}” overlaps “{B.text_frame.text[:22]}”")
            for sh in texts:
                if sh.top + sh.height > H:
                    faults.append(f"slide {i}: “{sh.text_frame.text[:22]}” runs off the bottom")
                for q in pics:
                    if (sh.left + sh.width > q.left and sh.top < q.top + q.height
                            and sh.top + sh.height > q.top):
                        faults.append(f"slide {i}: “{sh.text_frame.text[:22]}” runs under the picture")
        return faults

    def save(self, path: Path) -> Path:
        faults = self.check()
        if faults:
            raise SystemExit("layout is broken:\n  " + "\n  ".join(faults))
        self.prs.save(str(path))
        return path


# --- The deck -----------------------------------------------------------------
# Fifteen slides in four movements: the setup, the product, the business, the close.
# Two slides from the old seventeen are gone — "trust" repeated the money path, and
# "every way to grow in one place" repeated the product.
#
# Slides carrying an argument in columns or rows get no picture. Five rows of
# numbers plus a photograph is how the last deck ended up with a crumpled banknote
# illustrating a table about competitors.

def build() -> Deck:
    d = Deck()

    # -- setup ----------------------------------------------------------------
    d.title(
        [("where does your money ", False), ("sleep?", True)],
        "a non-custodial savings app on solana. hold dollars, earn real yield, "
        "own global assets — no bank, no broker, no crypto.",
        "eyes.png",
        "The eyes through torn paper: money watching, and doing nothing. It is the "
        "question itself, not decoration.",
    )

    d.statement(
        "the problem",
        [("inflation ", False), ("eats", True), (" your savings.", False)],
        "save in your local currency and you lose value every year. holding dollars "
        "that actually earn is the hard part.",
        False, "dripping-dollar.png",
        "Value visibly draining out of a note — loss over time, not 'money' in general.",
    )

    d.columns(
        "today's options",
        [("every door is ", False), ("half shut.", True)],
        [
            ("banks & neobanks", "4–5% at best, and the best usually needs a us bank account. everyone else gets a local-currency rate that inflation eats."),
            ("crypto wallets", "0%. your dollars sit still, and you carry the seed phrase, the gas and the scams yourself."),
            ("brokers", "treasuries, stocks, gold — real assets, gated by geography, paperwork and market hours."),
        ],
        False, None,
        "No picture: three columns already carry the argument.",
        footer="and the tools that do reach real yield are built for traders — seed phrases, gas, jargon, scams.",
    )

    d.statement(
        "who it's for",
        [("for the people the system ", False), ("forgets.", True)],
        "emerging-market savers, cross-border freelancers — anyone who wants dollars "
        "that grow, without becoming a crypto trader.",
        True, "crowd-hats.png",
        "One figure in violet picked out of the crowd. That single person is the slide.",
    )

    # -- product --------------------------------------------------------------
    d.statement(
        "the solution",
        [("a dollar account that ", False), ("actually earns.", True)],
        "one non-custodial account: hold dollars, earn yield, own treasuries, stocks "
        "and gold. email sign-in, apple pay, withdraw anytime.",
        False, "sleeping-money.png",
        "The product's own metaphor, and the one slide where a soft image is right.",
    )

    d.columns(
        "the product",
        [("four things, ", False), ("one account.", True)],
        [
            ("earn", "dollars into curated yield sources — jupiter lend, ondo usdy (us treasuries), maple (institutional credit), onre. 5–12% apy."),
            ("own", "tokenized stocks and gold, held in your own wallet. real price exposure, on-chain p&l, no broker."),
            ("fund it", "apple pay, card, or any crypto you already hold. gas is paid for you — you never need to buy sol."),
            ("your pile", "every position in one view, with what it earned. withdraw any of it, any time, without asking us."),
        ],
        True, None,
        "WANTS A REAL SCREENSHOT of the app, replacing nothing — this is where an "
        "investor wants proof the thing exists. Drop one on the right when we have it.",
        footer="one tap each, no apps to juggle, nothing to learn about crypto.",
    )

    d.steps(
        "how it works",
        [("your money never ", False), ("passes through us.", True)],
        [
            "sign in with an email. a wallet is created for you — no seed phrase to write down.",
            "fund it with apple pay, a card, or crypto you already hold.",
            "the money moves straight from your wallet into an audited protocol. the position is yours.",
        ],
        False, None,
        "No stock photo: this is a factual claim about custody, and a mood image only "
        "weakens it. A diagram we draw ourselves would work — wallet to protocol, with "
        "OXAR beside the arrow and not on it.",
        footer="oxar ships no smart contract of its own — we sit on top of audited protocols. "
               "nothing to hack, no keys for us to lose, no withdrawal for us to approve.",
    )

    # -- business -------------------------------------------------------------
    d.stats(
        "market",
        [("the dollars are ", False), ("already here, asleep.", True)],
        "the money is already on-chain, already in dollars, and already sitting still.",
        [
            ("$295B", "tam", "on-chain dollars earning nothing — about 95% of a $312b stablecoin market, held across 150m+ addresses"),
            ("$15B", "sam", "the idle share of solana's $16.7b stablecoin supply — money we can reach without asking anyone to bridge"),
            ("$3M", "som · year one", "two basis points of the idle solana float — about 1,400 savers at the $2,080 an average stablecoin address holds"),
        ],
        True, "eye-through-clouds.jpg",
        "Something that means 'asleep', not 'finance'. A skyline meant nothing; the eye "
        "through the sky at least watches.",
        footer="sources: artemis and visa onchain analytics, q2 2026. yield-bearing share is 2–6% "
               "depending on whose definition, so 95% idle is the conservative read.",
    )

    d.rows(
        "business model",
        [("a quarter of a percent, ", False), ("named before you sign.", True)],
        [
            ("the model", "0.25% on a conversion — buying or selling anything that has to be swapped, where dollars are one side of the trade. putting dollars into a dollar product is not a conversion and costs nothing.", False),
            ("why not yield", "the position sits in the user's own wallet on a public market, and we are not in the flow when they leave. a performance fee needs custody, and custody is the promise we sell.", False),
            ("the price", "revolut charges 1.49% to convert plus a 1.5–2.5% spread; phantom 0.85% hidden inside the quote; metamask 0.875%. we take 0.25%, as its own line, before you sign.", False),
            ("what it earns", "revenue tracks turnover, not deposits — about $1 for every $400 swapped. money that sleeps earns us nothing, which is the honest cost of building for savers.", False),
        ],
        False, None,
        "No picture. Four rows of argument is a full slide.",
        footer="built, disclosed in the terms, and switched off behind two switches — neither of them set. "
               "today, before launch, we take 0%.",
    )

    d.rows(
        "competition",
        [("nobody covers ", False), ("both halves.", True)],
        [
            ("revolut", "2.33% on dollars for a standard account, up to 3.68% only on a paid tier — with capital at risk if the fund's nav falls. converting costs 1.49% plus a 1.5–2.5% spread.", False),
            ("us savings apps", "4–5%, the honest ceiling of a treasury rate — and they need a us bank account, which is the whole problem for everyone else.", False),
            ("crypto wallets", "0% on the dollars sitting in them. phantom charges 0.85% to convert and hides it in the quote; metamask 0.875%.", False),
            ("defi frontends", "the yield is real and the fee is often zero — but the interface is built for people who already speak the language.", False),
            ("oxar", "5–12%, real assets, nothing held by us, and 0.25% to convert — a sixth of revolut, a third of phantom, printed on screen before you sign.", True),
        ],
        False, None,
        "No picture. Five rows, and the numbers are the argument.",
        footer="the defensible part is not the yield — anyone can route to a protocol. "
               "it is who we serve, and that they trust us with the first dollar.",
    )

    d.stats(
        "traction",
        [("small numbers, ", False), ("honestly counted.", True)],
        None,
        [
            ("live", "on solana mainnet", "not a prototype — real money, from your own wallet"),
            ("99+", "waitlist signups", "16 of them arrived through someone else's referral link, with nothing spent on acquisition"),
            ("$75k", "intended deposits", "self-reported by the ten who named a figure — an intention, not a commitment. median $1,200"),
            ("36", "assets", "5 yield sources, 28 tokenized stocks, gold and silver"),
        ],
        False, None,
        "The app, or nothing. Four figures with their caveats is dense already.",
        footer="counted against the production database on 16 august 2026, not remembered. gasless deposits, "
               "apple-pay funding, a live portfolio, english and ukrainian — shipped in months, bootstrapped.",
    )

    d.statement(
        "why now",
        [("the timing is ", False), ("now.", True)],
        "tokenized assets crossed into the billions, crypto payroll is mainstream, "
        "card-to-crypto onramps finally work. the window is open 12–24 months before "
        "revolut and coinbase close it.",
        True, "crowd-world.png",
        "Many people turning at once — 'everyone is arriving', not 'the future'.",
    )

    # -- close ----------------------------------------------------------------
    d.rows(
        "roadmap",
        [("the product is built. ", False), ("now the people.", True)],
        [
            ("now", "live on mainnet in closed alpha: dollar yield, tokenized stocks, gold — end to end, from your own wallet. apple pay funding and gasless deposits both shipped.", False),
            ("next", "open the door. the product is built and the gate is an allowlist; what is missing is people, not features.", False),
            ("q4 2026", "assets from issuers jupiter cannot route today, and from a us broker-dealer whose liquidity beats most of our shelf. euro yield, once lend positions are priced in their own currency.", False),
            ("2027", "native ios and android. more currencies than the dollar. geographic expansion through local partners.", False),
        ],
        False, None,
        "No picture, or the child pointing ahead if the slide feels bare.",
    )

    d.columns(
        "the team",
        [("two founders, one ", False), ("live product.", True)],
        [
            ("daniel lohachov", "product and engineering. founder of the prior oxar iteration; in the solana ecosystem since 2023."),
            ("anna tarapatska", "operations, legal and partnerships."),
        ],
        False, "logo-chrome.png",
        "The chrome mark, large and to the right. The one slide where the brand can "
        "simply be itself.",
        footer="building from ukraine, bootstrapped — a live product on mainnet in months.",
    )

    d.rows(
        "the ask",
        [("funding the launch, ", False), ("not the runway.", True)],
        [
            ("not raising", "no traditional angel round. a crypto vc seed is planned for post-pmf, 12–18 months out.", False),
            ("raising instead", "grants, accelerators and hackathon prizes — $30k to carry the launch.", False),
            ("in motion", "solana foundation ecosystem grant, colosseum. partnerships with delora, privy and kora already live in the product.", False),
        ],
        False, None,
        "The hand-drawn study of the mark would suit the close — but crop out the "
        "handwritten '4-18%' and '$230B' first, because neither is a number this deck "
        "uses.",
        footer="oxar.app · daniel.l@oxar.app · @eternaki",
    )

    return d


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.home() / "Desktop" / "OXAR-pitch.pptx"
    deck = build()
    deck.save(out)
    print(f"{out}  ·  {len(deck.prs.slides._sldIdLst)} slides")
